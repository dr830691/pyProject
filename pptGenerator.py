from pptx import Presentation
from pptx.util import Inches
def pptSlideAdder(pptname,n):
    prs = Presentation(pptname)
    title_slide_layout = prs.slide_layouts[1]
    for i in range(1,n):
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Smaple Title "+str(i)
        subtitle.text = "Sample Subtitle"+str(i)
        top =Inches(2.5)
        left = Inches(1)
        height = Inches(4.5)
        width=Inches(3.5)
        image_path="./compositeImages/new"+str(i)+".png"
        pic = slide.shapes.add_picture(image_path, left, top, height=height,width=width)
        print(str(i)+" slide is added successfully")
    prs.save('output.pptx')
    
# pptSlideAdder("new.png")